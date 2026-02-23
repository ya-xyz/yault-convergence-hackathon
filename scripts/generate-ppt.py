#!/usr/bin/env python3
"""
Generate Yault Guardian Vault — Convergence Hackathon Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)  # Deep navy
BG_CARD      = RGBColor(0x15, 0x1C, 0x28)  # Card background
ACCENT_BLUE  = RGBColor(0x37, 0x5B, 0xD2)  # Chainlink-ish blue
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)  # Bright cyan accent
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)  # Green for success
ACCENT_GOLD  = RGBColor(0xFF, 0xD7, 0x00)  # Gold highlight
TEXT_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT    = RGBColor(0xB0, 0xBC, 0xCE)  # Muted text
TEXT_DIM      = RGBColor(0x6B, 0x7B, 0x93)  # Dimmer text
ORANGE        = RGBColor(0xFF, 0x6B, 0x35)  # Warning / highlight
PURPLE        = RGBColor(0x7C, 0x3A, 0xED)  # Purple accent


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_LIGHT, spacing=Pt(8), bullet_color=ACCENT_CYAN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# Accent line at top
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

# Title
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(13), Inches(1.2),
             "Yault Guardian Vault", font_size=52, color=TEXT_WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(12), Inches(1.0),
             "Self-Custodial Inheritance & Conditional Asset Release",
             font_size=28, color=ACCENT_CYAN)

add_text_box(slide, Inches(1.5), Inches(4.2), Inches(12), Inches(0.8),
             "with auditable yield sharing, powered by Chainlink CRE attestations",
             font_size=20, color=TEXT_LIGHT)

# Track badges
add_shape_bg(slide, Inches(1.5), Inches(5.8), Inches(3.2), Inches(0.6), ACCENT_BLUE, 0.15)
add_text_box(slide, Inches(1.5), Inches(5.83), Inches(3.2), Inches(0.6),
             "DeFi & Tokenization", font_size=16, color=TEXT_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5.0), Inches(5.8), Inches(2.8), Inches(0.6), PURPLE, 0.15)
add_text_box(slide, Inches(5.0), Inches(5.83), Inches(2.8), Inches(0.6),
             "Risk & Compliance", font_size=16, color=TEXT_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Bottom info
add_text_box(slide, Inches(1.5), Inches(7.5), Inches(10), Inches(0.5),
             "Convergence: A Chainlink Hackathon  ·  February 2026",
             font_size=14, color=TEXT_DIM)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
             "The Problem", font_size=36, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(13), Inches(0.7),
             "Crypto inheritance & conditional release systems face a trilemma:",
             font_size=20, color=TEXT_LIGHT)

# 3 problem cards
problems = [
    ("Custodial Trust", "Traditional solutions require handing\nkeys to a third party", "🔐"),
    ("Weak Enforcement", "Purely social or legal mechanisms\nhave no on-chain teeth", "⚖️"),
    ("No Audit Trail", "Release decisions are opaque and\nhard to verify after the fact", "🔍"),
]

for i, (title, desc, icon) in enumerate(problems):
    x = Inches(1.2 + i * 4.6)
    y = Inches(2.8)
    w = Inches(4.2)
    h = Inches(3.5)
    add_shape_bg(slide, x, y, w, h, BG_CARD, 0.05)

    # Red accent line on top of card
    add_shape_bg(slide, x, y, w, Inches(0.06), ORANGE)

    add_text_box(slide, x + Inches(0.4), y + Inches(0.5), w - Inches(0.8), Inches(0.5),
                 icon, font_size=36, color=TEXT_WHITE)
    add_text_box(slide, x + Inches(0.4), y + Inches(1.2), w - Inches(0.8), Inches(0.5),
                 title, font_size=22, color=TEXT_WHITE, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(1.9), w - Inches(0.8), Inches(1.5),
                 desc, font_size=16, color=TEXT_LIGHT)

add_text_box(slide, Inches(1.2), Inches(7.2), Inches(13), Inches(0.6),
             "Result: $68B+ in crypto is estimated to be inaccessible due to lost keys, death, or custody failures.",
             font_size=16, color=ORANGE, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Our Approach
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
             "Our Approach: Zero-Custody Release", font_size=36, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(13), Inches(0.7),
             "No single party — not the platform, not the authority, not the recipient — can unilaterally access assets.",
             font_size=18, color=ACCENT_CYAN)

# Architecture layers table
layers = [
    ("Asset Custody", "Owner's wallet — keys never leave the client", ACCENT_GREEN),
    ("Key Protection", "AES-256-GCM-SIV + Argon2id (ACE-GF framework)", ACCENT_CYAN),
    ("Conditional Release", "Shamir Secret Sharing — authority holds a share, not a key", ACCENT_BLUE),
    ("Timelock Fallback", "drand BLS threshold IBE (tlock)", PURPLE),
    ("Attestation Trigger", "Chainlink CRE workflow → on-chain ReleaseAttestation", ACCENT_GOLD),
    ("Yield Management", "ERC-4626 vault with auditable user/platform/authority split", ACCENT_GREEN),
    ("Permanent Storage", "Arweave for encrypted release artifacts", ACCENT_CYAN),
]

for i, (layer, mechanism, accent) in enumerate(layers):
    y = Inches(2.5 + i * 0.72)

    # Accent bar
    add_shape_bg(slide, Inches(1.2), y, Inches(0.08), Inches(0.55), accent)

    add_text_box(slide, Inches(1.6), y + Inches(0.05), Inches(3.5), Inches(0.5),
                 layer, font_size=16, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(5.5), y + Inches(0.05), Inches(9.5), Inches(0.5),
                 mechanism, font_size=15, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — How Chainlink CRE Is Used
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(12), Inches(0.8),
             "How Chainlink CRE Is Used", font_size=36, color=TEXT_WHITE, bold=True)

# 3 integration boxes
integrations = [
    (
        "1. CRE Workflow",
        "oracle/workflow/src/main.ts",
        [
            "HTTP + Cron dual trigger mode",
            "Polls platform API for pending attestation requests",
            "Retry logic (3x exponential backoff)",
            "Encodes & submits via CRE EVM Write",
        ],
        ACCENT_BLUE
    ),
    (
        "2. On-Chain Attestation",
        "contracts/src/ReleaseAttestation.sol",
        [
            "Immutable attestation record",
            "Oracle source (source=0) + Fallback (source=1)",
            "keccak256 wallet ID for privacy",
            "Events emitted for platform consumption",
        ],
        ACCENT_CYAN
    ),
    (
        "3. Trigger Pipeline",
        "server/api/trigger/oracle.js",
        [
            "Reads on-chain attestation events",
            "Creates cooldown-gated triggers",
            "Separates event source from policy execution",
            "Authority verification before release",
        ],
        ACCENT_GREEN
    ),
]

for i, (title, file_path, points, accent) in enumerate(integrations):
    x = Inches(1.0 + i * 4.8)
    y = Inches(1.8)
    w = Inches(4.5)
    h = Inches(5.8)
    add_shape_bg(slide, x, y, w, h, BG_CARD, 0.04)
    add_shape_bg(slide, x, y, w, Inches(0.06), accent)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
                 title, font_size=20, color=TEXT_WHITE, bold=True)

    # File path
    add_shape_bg(slide, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), Inches(0.45), RGBColor(0x0A, 0x0E, 0x14), 0.08)
    add_text_box(slide, x + Inches(0.5), y + Inches(1.05), w - Inches(1.0), Inches(0.4),
                 file_path, font_size=11, color=ACCENT_CYAN, font_name='Consolas')

    for j, point in enumerate(points):
        add_text_box(slide, x + Inches(0.5), y + Inches(1.7 + j * 0.85), w - Inches(0.8), Inches(0.8),
                     f"• {point}", font_size=14, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — End-to-End Flow
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
             "End-to-End Release Flow", font_size=36, color=TEXT_WHITE, bold=True)

steps = [
    ("1", "Owner", "Configure recipient\npaths & deposit\nto vault", ACCENT_GREEN),
    ("2", "Chainlink CRE", "Oracle workflow\nwrites attestation\non-chain", ACCENT_BLUE),
    ("3", "Platform", "Read attestation\nevent → create\ncooldown trigger", ACCENT_CYAN),
    ("4", "Authority", "Verify conditions\n& submit release\nfactors (Shamir)", PURPLE),
    ("5", "Recipient", "Combine factors\n& claim released\nassets", ACCENT_GOLD),
    ("6", "Audit", "Yield split\nrecords visible\nfor all parties", ACCENT_GREEN),
]

for i, (num, actor, desc, accent) in enumerate(steps):
    x = Inches(0.6 + i * 2.5)
    y = Inches(2.2)
    w = Inches(2.2)
    h = Inches(4.5)

    # Card
    add_shape_bg(slide, x, y, w, h, BG_CARD, 0.04)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y + Inches(0.3), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.7), y + Inches(0.35), Inches(0.7), Inches(0.6),
                 num, font_size=22, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Actor name
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), w - Inches(0.4), Inches(0.5),
                 actor, font_size=16, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), w - Inches(0.4), Inches(2.2),
                 desc, font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(steps) - 1:
        add_text_box(slide, x + w - Inches(0.1), y + Inches(1.9), Inches(0.5), Inches(0.5),
                     "→", font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Key insight box
add_shape_bg(slide, Inches(1.2), Inches(7.2), Inches(13.5), Inches(0.8), BG_CARD, 0.05)
add_shape_bg(slide, Inches(1.2), Inches(7.2), Inches(0.08), Inches(0.8), ACCENT_GOLD)
add_text_box(slide, Inches(1.6), Inches(7.3), Inches(13), Inches(0.6),
             "Key: No single party can unilaterally access or move assets at any step.",
             font_size=16, color=ACCENT_GOLD, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — Tech Stack & Smart Contracts
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
             "Smart Contracts & Tech Stack", font_size=36, color=TEXT_WHITE, bold=True)

# Left: Contracts
add_text_box(slide, Inches(1.2), Inches(1.6), Inches(6), Inches(0.5),
             "On-Chain Contracts (Solidity 0.8.28)", font_size=20, color=ACCENT_CYAN, bold=True)

contracts = [
    ("YaultVault.sol", "ERC-4626 yield vault with 3-way revenue split\n(User 75% · Platform 20% · Authority 5%)"),
    ("YaultVaultFactory.sol", "Factory pattern for per-asset vault deployment"),
    ("ReleaseAttestation.sol", "Immutable on-chain attestation record\n(Oracle source + Fallback submitter)"),
    ("YaultPathClaim.sol", "Path-based asset claim with hash verification"),
    ("VaultShareEscrow.sol", "Escrow for ERC-4626 shares pending claim"),
]

for i, (name, desc) in enumerate(contracts):
    y = Inches(2.2 + i * 1.1)
    add_shape_bg(slide, Inches(1.2), y, Inches(7.0), Inches(0.95), BG_CARD, 0.04)
    add_text_box(slide, Inches(1.5), y + Inches(0.08), Inches(6.5), Inches(0.35),
                 name, font_size=14, color=ACCENT_GREEN, bold=True, font_name='Consolas')
    add_text_box(slide, Inches(1.5), y + Inches(0.4), Inches(6.5), Inches(0.55),
                 desc, font_size=12, color=TEXT_LIGHT)

# Right: Tech stack
add_text_box(slide, Inches(8.8), Inches(1.6), Inches(6), Inches(0.5),
             "Technology Stack", font_size=20, color=ACCENT_CYAN, bold=True)

stack = [
    ("Smart Contracts", "Solidity · Foundry · OpenZeppelin"),
    ("Backend", "Node.js · Express · sql.js"),
    ("Cryptography", "Rust/WASM · X25519 · AES-GCM-SIV · Shamir"),
    ("Oracle", "Chainlink CRE (TypeScript)"),
    ("Storage", "Arweave · AO"),
    ("Timelock", "drand network · tlock-js"),
    ("Frontend", "Vanilla JS · Web3 wallet connect"),
    ("Testing", "Foundry · Jest · Cargo"),
    ("Testnet", "Ethereum Sepolia · Tenderly"),
]

for i, (component, tech) in enumerate(stack):
    y = Inches(2.2 + i * 0.62)
    add_text_box(slide, Inches(9.0), y, Inches(2.8), Inches(0.5),
                 component, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(11.8), y, Inches(3.5), Inches(0.5),
                 tech, font_size=13, color=TEXT_LIGHT)

# Deployed badge
add_shape_bg(slide, Inches(8.8), Inches(7.6), Inches(6.0), Inches(0.5), ACCENT_GREEN, 0.15)
add_text_box(slide, Inches(8.8), Inches(7.63), Inches(6.0), Inches(0.45),
             "✓  All contracts deployed & verified on Sepolia + Tenderly", font_size=14,
             color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — Cryptographic Security Model
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(12), Inches(0.8),
             "Cryptographic Security Model", font_size=36, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(13), Inches(0.7),
             "Authorities are equivalent to drand BLS signing nodes — they hold a protocol share with zero standalone capability.",
             font_size=17, color=ACCENT_CYAN)

# Security layers
sec_layers = [
    ("ACE-GF Framework", "AES-256-GCM-SIV encryption with Argon2id key derivation.\nAdmin factors encrypted per-authority with unique salt.", ACCENT_BLUE),
    ("Shamir Secret Sharing", "Key split into N shares with threshold T.\nNo single share reveals anything about the original key.", PURPLE),
    ("End-to-End Encryption", "X25519 key agreement + ChaCha20-Poly1305 for\nauthority ↔ recipient communication. Platform sees nothing.", ACCENT_GREEN),
    ("Timelock Fallback", "drand network BLS-IBE (tlock): if authority is unavailable,\nshares auto-decrypt after pre-set time. No trusted party needed.", ACCENT_GOLD),
    ("On-Chain Attestation", "Chainlink oracle writes immutable attestation.\nCooldown period prevents premature release.", ACCENT_CYAN),
]

for i, (title, desc, accent) in enumerate(sec_layers):
    y = Inches(2.5 + i * 1.2)
    add_shape_bg(slide, Inches(1.2), y, Inches(13.5), Inches(1.05), BG_CARD, 0.04)
    add_shape_bg(slide, Inches(1.2), y, Inches(0.08), Inches(1.05), accent)
    add_text_box(slide, Inches(1.6), y + Inches(0.1), Inches(4.0), Inches(0.4),
                 title, font_size=17, color=accent, bold=True)
    add_text_box(slide, Inches(5.8), y + Inches(0.1), Inches(8.5), Inches(0.9),
                 desc, font_size=14, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — Revenue Model & Yield Split
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
             "ERC-4626 Yield Vault & Revenue Split", font_size=36, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(13), Inches(0.7),
             "Assets in the vault earn yield while waiting for release — split transparently among all participants.",
             font_size=18, color=TEXT_LIGHT)

# Revenue split visualization
splits = [
    ("User (Vault Depositor)", "75%", Inches(8.0), ACCENT_GREEN),
    ("Platform (Yault)", "20%", Inches(2.1), ACCENT_BLUE),
    ("Authority (Witness)", "5%", Inches(0.55), PURPLE),
]

bar_y = Inches(3.2)
bar_x = Inches(1.5)
total_w = Inches(10.65)
current_x = bar_x

for label, pct, width, color in splits:
    add_shape_bg(slide, current_x, bar_y, width, Inches(1.0), color, 0.03)
    add_text_box(slide, current_x, bar_y + Inches(0.15), width, Inches(0.7),
                 f"{pct}", font_size=28, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, current_x, bar_y + Inches(1.15), width, Inches(0.5),
                 label, font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    current_x = current_x + width + Inches(0.0)

# Key features
features = [
    "ERC-4626 compliant — composable with any DeFi yield strategy",
    "Per-asset vaults via factory pattern — USDC, WETH, etc.",
    "Yield accrues during cooldown period — no idle capital",
    "Authority incentivized to participate via 5% share",
    "All splits on-chain and auditable in vault accounting",
]

add_text_box(slide, Inches(1.2), Inches(5.2), Inches(6), Inches(0.5),
             "Key Features", font_size=20, color=ACCENT_CYAN, bold=True)

for i, feat in enumerate(features):
    add_text_box(slide, Inches(1.5), Inches(5.8 + i * 0.5), Inches(13), Inches(0.45),
                 f"✓  {feat}", font_size=15, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — Landscape Comparison
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.7),
             "Landscape: Why Existing Solutions Fall Short", font_size=34, color=TEXT_WHITE, bold=True)

# Column headers
col_headers = ["Approach", "Self-\nCustody", "Yield", "Auto\nTrigger", "Low Trust\nDep.", "Revoc-\nable", "Main Weakness"]
col_x =       [0.3,        4.8,         6.2,     7.6,     9.0,         10.4,        11.8]
col_w =       [4.3,        1.2,         1.2,     1.2,     1.2,         1.2,         3.8]

for j, header in enumerate(col_headers):
    x = Inches(col_x[j])
    w = Inches(col_w[j])
    add_shape_bg(slide, x, Inches(1.3), w, Inches(0.65), ACCENT_BLUE, 0.03)
    add_text_box(slide, x + Inches(0.05), Inches(1.3), w - Inches(0.1), Inches(0.65),
                 header, font_size=10, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Status symbols
S_YES = "✓"      # green
S_NO  = "—"      # dim
S_PART = "~"     # yellow/partial

# Data: (approach, custody, yield, trigger, low_dep, revocable, weakness)
comp_rows = [
    ("Centralized Custodial\nInheritance",   S_NO,   S_PART, S_NO,   S_NO,   S_NO,   "Platform risk; policy changes;\nnot self-custody"),
    ("Legal / Trust\nCustody",               S_NO,   S_NO,   S_NO,   S_NO,   S_PART, "Single point of trust;\nkey exposure; slow; expensive"),
    ("Multi-Sig\nInheritance",               S_YES,  S_PART, S_NO,   S_PART, S_PART, "Coordination burden;\nsigner collusion; no yield mgmt"),
    ("MPC\nCustody",                         S_PART, S_YES,  S_NO,   S_NO,   S_NO,   "High cost; service disruption;\nstill centralized nodes"),
    ("Dead Man's\nSwitch",                   S_YES,  S_NO,   S_YES,  S_PART, S_PART, "Time ≠ death; irreversible\nmisfire; can't manage DeFi"),
    ("Social\nRecovery",                     S_YES,  S_YES,  S_NO,   S_PART, S_YES,  "Collusion risk; not true\ninheritance; no yield handoff"),
    ("DeFi Vaults\n(non-inheritance)",       S_YES,  S_YES,  S_NO,   S_YES,  S_YES,  "No inheritance logic;\nliquidation risk"),
]

# Yault row (highlighted)
yault_row = ("Yault Guardian\nVault", S_YES, S_YES, S_YES, S_YES, S_YES, "Hackathon build;\nrequires production audit")

def status_color(s):
    if s == S_YES:
        return ACCENT_GREEN
    elif s == S_PART:
        return ACCENT_GOLD
    else:
        return RGBColor(0x4A, 0x55, 0x68)  # dim for "no"

row_h = 0.72
start_y = 2.05

for i, row_data in enumerate(comp_rows):
    y = Inches(start_y + i * row_h)
    row_bg = BG_CARD if i % 2 == 0 else BG_DARK

    # Row background
    add_shape_bg(slide, Inches(0.3), y, Inches(15.3), Inches(row_h - 0.02), row_bg, 0.02)

    for j, val in enumerate(row_data):
        x = Inches(col_x[j])
        w = Inches(col_w[j])
        if j == 0:
            add_text_box(slide, x + Inches(0.08), y + Inches(0.02), w - Inches(0.1), Inches(row_h),
                         val, font_size=10, color=TEXT_WHITE, bold=False)
        elif j == len(row_data) - 1:  # weakness column
            add_text_box(slide, x + Inches(0.08), y + Inches(0.02), w - Inches(0.1), Inches(row_h),
                         val, font_size=9, color=TEXT_DIM)
        else:
            c = status_color(val)
            add_text_box(slide, x, y + Inches(0.08), w, Inches(row_h - 0.1),
                         val, font_size=16, color=c, bold=True, alignment=PP_ALIGN.CENTER)

# Yault row — highlighted
y_yault = Inches(start_y + len(comp_rows) * row_h + 0.08)
add_shape_bg(slide, Inches(0.3), y_yault, Inches(15.3), Inches(row_h + 0.05), ACCENT_BLUE, 0.02)

for j, val in enumerate(yault_row):
    x = Inches(col_x[j])
    w = Inches(col_w[j])
    if j == 0:
        add_text_box(slide, x + Inches(0.08), y_yault + Inches(0.04), w - Inches(0.1), Inches(row_h),
                     val, font_size=11, color=TEXT_WHITE, bold=True)
    elif j == len(yault_row) - 1:
        add_text_box(slide, x + Inches(0.08), y_yault + Inches(0.04), w - Inches(0.1), Inches(row_h),
                     val, font_size=9, color=TEXT_LIGHT)
    else:
        add_text_box(slide, x, y_yault + Inches(0.1), w, Inches(row_h - 0.1),
                     val, font_size=16, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom legend
add_text_box(slide, Inches(0.8), Inches(8.2), Inches(14), Inches(0.35),
             "✓ = fully supported      ~ = partial / depends      — = not supported",
             font_size=11, color=TEXT_DIM)

add_text_box(slide, Inches(0.8), Inches(8.5), Inches(14), Inches(0.35),
             "Yault is the only approach combining self-custody, native yield, oracle-driven triggers, low trust dependency, and revocability.",
             font_size=12, color=ACCENT_CYAN, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — Testing & Quality
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
             "Testing & Quality Assurance", font_size=36, color=TEXT_WHITE, bold=True)

# Test stats
test_cats = [
    ("14", "Integration\nTests", "API + flow tests\n(Jest + Supertest)", ACCENT_BLUE),
    ("11", "Unit\nTests", "Module-level tests\n(Jest)", ACCENT_CYAN),
    ("5", "Contract\nTests", "Solidity tests\n(Foundry/Forge)", ACCENT_GREEN),
    ("2", "WASM\nTests", "Crypto primitive tests\n(Cargo)", PURPLE),
]

for i, (count, label, desc, accent) in enumerate(test_cats):
    x = Inches(1.0 + i * 3.7)
    y = Inches(1.9)
    w = Inches(3.4)
    h = Inches(3.0)
    add_shape_bg(slide, x, y, w, h, BG_CARD, 0.04)
    add_text_box(slide, x, y + Inches(0.3), w, Inches(0.8),
                 count, font_size=48, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.2), w, Inches(0.6),
                 label, font_size=16, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(2.0), w, Inches(0.8),
                 desc, font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Security review note
add_shape_bg(slide, Inches(1.2), Inches(5.5), Inches(13.5), Inches(2.5), BG_CARD, 0.04)
add_shape_bg(slide, Inches(1.2), Inches(5.5), Inches(0.08), Inches(2.5), ORANGE)

add_text_box(slide, Inches(1.6), Inches(5.7), Inches(12), Inches(0.5),
             "Internal Security Review Completed", font_size=18, color=ORANGE, bold=True)

sec_items = [
    "3 rounds of code review (P0 → P1 → P2 severity levels)",
    "All P0 critical issues identified and remediated",
    "Known limitations documented with mitigation strategies",
    "Production deployment requires third-party audit (planned post-hackathon)",
]
for i, item in enumerate(sec_items):
    add_text_box(slide, Inches(1.9), Inches(6.3 + i * 0.45), Inches(12), Inches(0.4),
                 f"•  {item}", font_size=14, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — Demo / Live Flow
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
             "Live Demo", font_size=36, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(13), Inches(0.7),
             "Full golden-path walkthrough: Owner → Oracle → Authority → Recipient → Audit",
             font_size=20, color=ACCENT_CYAN)

demo_steps = [
    ("Step 1", "Owner Portal", "Configure recipient path, set conditions,\ndeposit assets into ERC-4626 vault", ACCENT_GREEN),
    ("Step 2", "CRE Workflow", "Oracle attests conditions are met →\nwrite ReleaseAttestation on-chain", ACCENT_BLUE),
    ("Step 3", "Trigger Pipeline", "Platform reads attestation event →\ncreate trigger with cooldown period", ACCENT_CYAN),
    ("Step 4", "Authority Portal", "Verify conditions → submit Shamir\nrelease factors (authenticated)", PURPLE),
    ("Step 5", "Recipient Claim", "Combine factors → claim vault shares\nfrom escrow → withdraw assets", ACCENT_GOLD),
    ("Step 6", "Audit & Yield", "View revenue split records:\n75% user · 20% platform · 5% authority", ACCENT_GREEN),
]

for i, (step, title, desc, accent) in enumerate(demo_steps):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 5.0)
    y = Inches(2.5 + row * 3.2)
    w = Inches(4.7)
    h = Inches(2.8)

    add_shape_bg(slide, x, y, w, h, BG_CARD, 0.04)
    add_shape_bg(slide, x, y, w, Inches(0.06), accent)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(1.2), Inches(0.4),
                 step, font_size=12, color=TEXT_DIM, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.65), w - Inches(0.6), Inches(0.5),
                 title, font_size=18, color=accent, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.3), w - Inches(0.6), Inches(1.3),
                 desc, font_size=14, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 12 — Closing / Thank You
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(16), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(13), Inches(1.0),
             "Yault Guardian Vault", font_size=48, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(13), Inches(0.8),
             "Bringing trust-minimized, auditable asset release to crypto inheritance.",
             font_size=22, color=ACCENT_CYAN)

# Summary points
summary = [
    "Self-custodial — no party holds keys or unilateral power",
    "Chainlink CRE attestation — trusted, externalized trigger signal",
    "Auditable — every decision recorded on-chain + Arweave",
    "Yield-generating — assets earn while waiting for release",
]

for i, point in enumerate(summary):
    add_text_box(slide, Inches(1.8), Inches(4.0 + i * 0.6), Inches(12), Inches(0.5),
                 f"✦  {point}", font_size=18, color=TEXT_LIGHT)

# Links section
add_shape_bg(slide, Inches(1.5), Inches(6.6), Inches(13), Inches(1.5), BG_CARD, 0.05)

add_text_box(slide, Inches(2.0), Inches(6.8), Inches(5), Inches(0.4),
             "GitHub", font_size=14, color=TEXT_DIM, bold=True)
add_text_box(slide, Inches(4.0), Inches(6.8), Inches(8), Inches(0.4),
             "github.com/hayekw/yault-convergence-hackathon", font_size=14, color=ACCENT_CYAN)

add_text_box(slide, Inches(2.0), Inches(7.2), Inches(5), Inches(0.4),
             "Demo Video", font_size=14, color=TEXT_DIM, bold=True)
add_text_box(slide, Inches(4.0), Inches(7.2), Inches(8), Inches(0.4),
             "[link to be added]", font_size=14, color=TEXT_LIGHT)

add_text_box(slide, Inches(2.0), Inches(7.6), Inches(5), Inches(0.4),
             "Tenderly", font_size=14, color=TEXT_DIM, bold=True)
add_text_box(slide, Inches(4.0), Inches(7.6), Inches(8), Inches(0.4),
             "[link to be added]", font_size=14, color=TEXT_LIGHT)

# Bottom
add_text_box(slide, Inches(1.5), Inches(8.3), Inches(13), Inches(0.4),
             "Thank you — we welcome technical and security questions.",
             font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════
output_dir = "/Users/jwang/dev.y/yault-convergence-hackathon"
output_path = os.path.join(output_dir, "Yault_Guardian_Vault_Presentation.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
